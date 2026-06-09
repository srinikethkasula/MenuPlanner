"""
ui/ppt_preview.py — Live slide preview using actual template design.
Extracts background images from template.pptx and composites dish content on top.
"""
from __future__ import annotations
import io
from pathlib import Path
from typing import Optional

from PyQt6.QtWidgets import (
    QWidget, QVBoxLayout, QHBoxLayout, QLabel, QListWidget,
    QListWidgetItem, QPushButton, QFrame,
)
from PyQt6.QtCore import Qt, QThread, pyqtSignal, QSize
from PyQt6.QtGui import QPixmap, QImage, QFont, QIcon

try:
    from PIL import Image, ImageDraw, ImageFont
    PIL_OK = True
except ImportError:
    PIL_OK = False

from config import DAYS, SLOTS, PPT_SECTION_COLORS, RESOURCES_DIR
import database as db

# ── constants ─────────────────────────────────────────────────────────────────
_TEMPLATE = RESOURCES_DIR / "template.pptx"
_SW, _SH = 12192000, 6858000   # slide EMU dimensions
PW, PH   = 640, 360             # preview pixel dimensions

# Template colours (from template XML)
_C_NAME   = (51,  42,  37)    # #332A25 – dish name
_C_OLIVE  = (64,  94,   1)    # #405E01 – nutrition text
_C_WHITE  = (255, 255, 255)
_C_CREAM  = (253, 248, 243)   # background fallback

# Background cache: slide_idx → PIL Image (keyed by (idx, skip_tuple))
_BG: dict[tuple, Image.Image] = {}


def clear_bg_cache():
    """Call when the template file changes."""
    _BG.clear()


# ── EMU → pixel helpers ───────────────────────────────────────────────────────

def _px(emu): return int(emu * PW / _SW)
def _py(emu): return int(emu * PH / _SH)
def _pw(emu): return max(1, int(emu * PW / _SW))
def _ph(emu): return max(1, int(emu * PH / _SH))


# ── template background extraction ───────────────────────────────────────────

def _build_bg(slide_idx: int, skip: tuple = ()) -> "Image.Image":
    """Composite all picture-fill shapes from template slide onto a canvas."""
    canvas = Image.new('RGBA', (PW, PH), (*_C_CREAM, 255))
    if not PIL_OK or not _TEMPLATE.exists():
        return canvas
    try:
        from pptx import Presentation
        from pptx.oxml.ns import qn
        prs = Presentation(str(_TEMPLATE))
        slide = prs.slides[slide_idx]

        def _paste(shape, ox=0, oy=0):
            if shape.name in skip:
                return
            if shape.shape_type == 6:          # group — recurse
                try:
                    for s in shape.shapes:
                        _paste(s, ox + shape.left, oy + shape.top)
                except Exception:
                    pass
                return
            bl = shape._element.find('.//' + qn('a:blip'))
            if bl is None:
                return
            rId = bl.get(qn('r:embed'))
            try:
                blob = slide.part.rels[rId].target_part.blob
                src  = Image.open(io.BytesIO(blob)).convert('RGBA')
                x = _px(shape.left + ox);  y = _py(shape.top + oy)
                w = _pw(shape.width);       h = _ph(shape.height)
                src = src.resize((w, h), Image.LANCZOS)
                canvas.paste(src, (x, y), src)
            except Exception:
                pass

        for s in slide.shapes:
            _paste(s)
    except Exception:
        pass
    return canvas


import threading
_BG_LOCK = threading.Lock()

def _get_bg(slide_idx: int, skip: tuple = ()) -> "Image.Image":
    key = (slide_idx, skip)
    with _BG_LOCK:
        if key not in _BG:
            _BG[key] = _build_bg(slide_idx, skip)
        return _BG[key].copy()


def prewarm_cache():
    """Pre-load slides to memory cache."""
    try:
        _get_bg(0, skip=('Freeform 12',))
        _get_bg(1, skip=('Freeform 10', 'Freeform 11'))
        _get_bg(2)
    except Exception:
        pass


class PrewarmThread(QThread):
    def run(self):
        prewarm_cache()


# ── font helpers ─────────────────────────────────────────────────────────────

def _font(size: int, bold=False):
    win_fonts = Path("C:/Windows/Fonts")
    names = ["arialbd.ttf", "arial.ttf"] if bold else ["arial.ttf", "arialbd.ttf"]
    candidates = [str(win_fonts / n) for n in names] + list(names) + ["/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"]
    for p in candidates:
        try:
            return ImageFont.truetype(p, size)
        except Exception:
            pass
    return ImageFont.load_default()


def _pil_to_qpix(img: "Image.Image") -> QPixmap:
    if img.mode != "RGBA":
        img = img.convert("RGBA")
    data = img.tobytes("raw", "RGBA")
    qi = QImage(data, img.width, img.height, img.width * 4,
                QImage.Format.Format_RGBA8888)
    return QPixmap.fromImage(qi)


# ── slide renderers ───────────────────────────────────────────────────────────

# Dish slide layout (EMU → px, pre-computed)
_GRP9_X, _GRP9_Y = _px(817489),  _py(772106)
_GRP9_W, _GRP9_H = _pw(5044482), _ph(4987189)
_NAME_X,  _NAME_Y  = _px(5785141), _py(1701193)
_SEC_X,   _SEC_Y   = _px(7768057), _py(468797)
_NFACT_X, _NFACT_Y = _px(8149062), _py(2790262)
_LBL_X  = _px(6818332)
_VAL_X  = _px(9205906)
_ROW_TOPS = [
    _py(3550765),   # CALORIES
    _py(4222621),   # PROTEIN
    _py(4896571),   # FATS
    _py(5608387),   # FIBRE
    _py(6280000),   # CARBS (added below FIBRE)
]
_NUTRITION_ROWS = ['Calories', 'Protein', 'Fats', 'Fibre', 'Carbs']
_DISH_KEYS      = ['calories', 'protein', 'fat', 'fiber', 'carbs']
_UNITS          = ['KCAL', 'g', 'g', 'g', 'g']


def render_good_morning(welcome_text: str) -> QPixmap:
    bg = _get_bg(0, skip=('Freeform 12',))
    d  = ImageDraw.Draw(bg)
    # "Good Morning" in template position
    d.text((_px(959363), _py(1174590)), "Good Morning",
           font=_font(28, bold=True), fill=_C_NAME)
    # Welcome description
    for i, line in enumerate(welcome_text.split('\n')[:3]):
        d.text((_px(457200), _py(3695017) + i * 22), line,
               font=_font(11), fill=_C_NAME)
    return _pil_to_qpix(bg)


def _draw_wrap_text(d, text: str, x: int, y: int, font, max_w: int, fill):
    words = text.split()
    lines = []
    curr = []
    for w in words:
        test = " ".join(curr + [w])
        try:
            bbox = font.getbbox(test)
            width = bbox[2] - bbox[0]
        except Exception:
            try:
                width = d.textsize(test, font=font)[0]
            except Exception:
                width = len(test) * (font.size * 0.6)
        if width <= max_w:
            curr.append(w)
        else:
            if curr:
                lines.append(" ".join(curr))
                curr = [w]
            else:
                lines.append(w)
                curr = []
    if curr:
        lines.append(" ".join(curr))
        
    cy = y
    line_h = font.size + 2
    for line in lines:
        d.text((x, cy), line, font=font, fill=fill)
        cy += line_h


def render_dish_slide(dish: dict, category: str) -> QPixmap:
    # Background without dish photo area
    bg = _get_bg(1, skip=('Freeform 10', 'Freeform 11'))
    d  = ImageDraw.Draw(bg)

    # Dish photo — circular crop at Group 9 bounds
    img_path = dish.get('image_path', '')
    from database import resolve_image_path
    resolved = resolve_image_path(img_path)
    if resolved and Path(resolved).exists():
        try:
            photo = Image.open(resolved).convert('RGBA')
            photo = photo.resize((_GRP9_W, _GRP9_H), Image.LANCZOS)
            mask  = Image.new('L', (_GRP9_W, _GRP9_H), 0)
            ImageDraw.Draw(mask).ellipse([0, 0, _GRP9_W - 1, _GRP9_H - 1], fill=255)
            bg.paste(photo, (_GRP9_X, _GRP9_Y), mask)
        except Exception:
            pass

    # Section label (white, bold)
    d.text((_SEC_X, _SEC_Y + 8), category.upper(),
           font=_font(16, bold=True), fill=_C_WHITE)

    # Dish name
    name = dish.get('name', '').upper()
    _draw_wrap_text(d, name, _NAME_X, _NAME_Y, _font(14, bold=True), _NFACT_X - _NAME_X - 10, _C_NAME)

    # "NUTRITION FACTS" header
    d.text((_NFACT_X, _NFACT_Y), "NUTRITION FACTS",
           font=_font(10, bold=True), fill=_C_OLIVE)

    # Nutrition rows (5 rows)
    for i, (lbl, key, unit) in enumerate(
            zip(_NUTRITION_ROWS, _DISH_KEYS, _UNITS)):
        if i >= len(_ROW_TOPS):
            break
        ry = _ROW_TOPS[i]
        v  = dish.get(key, '—')
        val_str = f'{v} {unit}' if v not in ('—', '', None) else '—'
        d.text((_LBL_X, ry), lbl,     font=_font(9), fill=_C_OLIVE)
        d.text((_VAL_X, ry), val_str, font=_font(9, bold=True), fill=_C_OLIVE)

    return _pil_to_qpix(bg)


def render_closing() -> QPixmap:
    bg = _get_bg(2)
    d  = ImageDraw.Draw(bg)
    co = db.get_setting("company_name") or "Chathradhari Caterers"
    d.text((_px(2155419), _py(689568)), "Meet the Team",
           font=_font(22, bold=True), fill=_C_NAME)
    d.text((_px(2155419), _py(1200000)), co,
           font=_font(14), fill=_C_OLIVE)
    return _pil_to_qpix(bg)


# ── render thread ─────────────────────────────────────────────────────────────

class RenderThread(QThread):
    slide_ready = pyqtSignal(int, QPixmap, str)

    def __init__(self, slides: list):
        super().__init__()
        self.slides   = slides
        self._welcome = db.get_setting("welcome_text") or "Welcome to Oyster Cafeteria"
        self._cancelled = False

    def cancel(self):
        """Cooperatively request thread loop cancellation."""
        self._cancelled = True

    def run(self):
        for i, (kind, dish, label) in enumerate(self.slides):
            if self._cancelled:
                break
            try:
                if kind == "good_morning":
                    pix = render_good_morning(self._welcome)
                elif kind == "dish":
                    pix = render_dish_slide(dish, dish.get("_category", "Lunch"))
                else:
                    pix = render_closing()
            except Exception:
                pix = _pil_to_qpix(Image.new('RGBA', (PW, PH), (*_C_CREAM, 255)))
            
            if self._cancelled:
                break
            self.slide_ready.emit(i, pix, label)


# ── PPTPreview widget ─────────────────────────────────────────────────────────

class PPTPreview(QWidget):
    export_requested = pyqtSignal(str)

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setFixedWidth(360)
        self.setStyleSheet(
            "PPTPreview { background: #121212; border-left: 1px solid #252525; }"
        )
        self._current_day  = DAYS[0]
        self._day_data:    dict = {}
        self._dishes_db:   dict = {}
        self._thread:      Optional[RenderThread] = None
        self._slide_items: list[QListWidgetItem]  = []
        self._setup_ui()
        
        # Pre-warm templates in the background
        self._prewarm_thread = PrewarmThread(self)
        self._prewarm_thread.start()

    def _setup_ui(self):
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(0)

        # Title bar
        tb = QFrame()
        tb.setFixedHeight(54)
        tb.setStyleSheet(
            "QFrame { background: #1a1a1a; border-bottom: 1px solid #252525; }"
        )
        trow = QHBoxLayout(tb)
        trow.setContentsMargins(14, 0, 14, 0)
        tl = QLabel("Slide Preview")
        tl.setFont(QFont("Segoe UI", 12, QFont.Weight.Bold))
        tl.setStyleSheet("color: #e0e0e0; background: transparent;")
        trow.addWidget(tl)
        trow.addStretch()
        self._count_lbl = QLabel("")
        self._count_lbl.setStyleSheet(
            "color: #4a9eff; font-size: 10px; background: transparent;"
        )
        trow.addWidget(self._count_lbl)
        layout.addWidget(tb)

        # Day tabs
        tab_frame = QFrame()
        tab_frame.setFixedHeight(42)
        tab_frame.setStyleSheet(
            "QFrame { background: #1a1a1a; border-bottom: 1px solid #252525; }"
        )
        tab_row = QHBoxLayout(tab_frame)
        tab_row.setContentsMargins(8, 5, 8, 5)
        tab_row.setSpacing(4)
        self._day_btns: dict[str, QPushButton] = {}
        for day in DAYS:
            b = QPushButton(day[:3])
            b.setCheckable(True)
            b.setChecked(day == DAYS[0])
            b.setFixedSize(50, 28)
            b.setCursor(Qt.CursorShape.PointingHandCursor)
            b.clicked.connect(lambda _, d=day: self._select_day(d))
            b.setStyleSheet("""
                QPushButton {
                    background: #252525; border: 1px solid #333;
                    border-radius: 6px; color: #888;
                    font-size: 10px; font-weight: bold;
                }
                QPushButton:hover  { background: #303030; color: #ccc; }
                QPushButton:checked {
                    background: #1a3a6e; border-color: #4a9eff; color: #4a9eff;
                }
            """)
            tab_row.addWidget(b)
            self._day_btns[day] = b
        tab_row.addStretch()
        layout.addWidget(tab_frame)

        # Slide list
        self.slide_list = QListWidget()
        self.slide_list.setIconSize(QSize(320, 180))
        self.slide_list.setSpacing(5)
        self.slide_list.setStyleSheet("""
            QListWidget {
                background: #121212; border: none; padding: 6px;
            }
            QListWidget::item {
                background: #1a1a1a; border: 1px solid #252525;
                border-radius: 8px; color: #888;
                font-size: 10px; padding: 4px; margin: 2px 0;
            }
            QListWidget::item:selected {
                background: #1a3a6e; border: 1px solid #4a9eff; color: #4a9eff;
            }
            QListWidget::item:hover { background: #1f1f1f; }
        """)
        layout.addWidget(self.slide_list, 1)

        if not PIL_OK:
            w = QLabel("pip install Pillow\nfor live preview")
            w.setAlignment(Qt.AlignmentFlag.AlignCenter)
            w.setStyleSheet(
                "color: #ff7043; font-size: 10px; padding: 10px; background: transparent;"
            )
            layout.addWidget(w)

    def _select_day(self, day: str):
        self._current_day = day
        for d, b in self._day_btns.items():
            b.setChecked(d == day)
        self._render()

    def update_data(self, day_data: dict, dishes_db: dict):
        self._day_data  = day_data
        self._dishes_db = dishes_db
        self._render()

    def _build_slides(self) -> list:
        slides    = [("good_morning", {}, "Good Morning")]
        day_slots = self._day_data.get(self._current_day, {})
        order     = ["Breakfast", "Juice", "Lunch", "Salad", "Soup",
                     "Dessert", "Evening Snacks", "Fruit Lunch"]
        buckets: dict[str, list] = {s: [] for s in order}
        for slot_id, _, slot_sec, ppt_sec, _ in SLOTS:
            names_str = day_slots.get(slot_id, "").strip()
            if names_str:
                # Split by '+' to handle multiple dishes in a single slot
                names = [n.strip() for n in names_str.split("+") if n.strip()]
                for name in names:
                    info = dict(self._dishes_db.get(name, {}))
                    info.setdefault("name", name)
                    
                    db_cat = info.get("category", "")
                    
                    # Apply dynamic section-wise slide labeling rules
                    if slot_sec == "FRUIT":
                        slide_sec = "Fruit Lunch"
                    elif db_cat in ("Dessert", "Sweet"):
                        slide_sec = "Dessert"
                    elif db_cat == "Juice":
                        slide_sec = "Juice"
                    elif slot_sec == "LUNCH":
                        slide_sec = "Lunch"
                    elif slot_id in ("MORN1", "MORN2", "MORN3"):
                        slide_sec = "Breakfast"
                    elif slot_sec == "ADDITIONAL":
                        slide_sec = "Breakfast"
                    else:
                        slide_sec = ppt_sec
                        
                    info["_category"] = slide_sec
                    if slide_sec in buckets:
                        buckets[slide_sec].append(info)
        for sec in order:
            for dish in buckets[sec]:
                slides.append(("dish", dish, dish.get("name", "Dish")))
        slides.append(("closing", {}, "Team Photo / Logo"))
        return slides

    def _render(self):
        if not PIL_OK:
            return
        if self._thread and self._thread.isRunning():
            try:
                self._thread.slide_ready.disconnect()
            except TypeError:
                pass  # already disconnected
            self._thread.cancel()

        slides = self._build_slides()
        self.slide_list.clear()
        self._slide_items = []
        for i, (_, _, label) in enumerate(slides):
            item = QListWidgetItem(f"  [{i + 1}]  {label}")
            item.setSizeHint(QSize(335, 46))
            self.slide_list.addItem(item)
            self._slide_items.append(item)
        self._count_lbl.setText(f"{len(slides)} slides")

        self._thread = RenderThread(slides)
        self._thread.slide_ready.connect(self._on_ready)
        self._thread.start()

    def _on_ready(self, idx: int, pix: QPixmap, label: str):
        if idx >= len(self._slide_items):
            return
        item   = self._slide_items[idx]
        scaled = pix.scaled(320, 180,
                            Qt.AspectRatioMode.KeepAspectRatio,
                            Qt.TransformationMode.SmoothTransformation)
        item.setIcon(QIcon(scaled))
        item.setText(f"  [{idx + 1}]  {label}")
        item.setSizeHint(QSize(335, 196))
        self.slide_list.setIconSize(QSize(320, 180))
