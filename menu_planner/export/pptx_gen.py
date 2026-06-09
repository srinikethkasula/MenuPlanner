"""
pptx_gen.py — Template-based PPT generation.
Loads resources/template.pptx, clones its 4 slides for each output slide,
then removes the original 4. All fonts, theme and design stay intact.
"""
from __future__ import annotations
from pathlib import Path
from datetime import date
import copy

from lxml import etree as _et
from pptx import Presentation
from pptx.oxml.ns import qn

from config import SLOTS, RESOURCES_DIR

TEMPLATE_PATH = RESOURCES_DIR / "template.pptx"

# Template slide indices (must stay fixed during build; deleted at the end)
_T_GM    = 0   # Good Morning
_T_DISH  = 1   # Dish slide
_T_TEAM  = 2   # Team photo
_T_LOGO  = 3   # Logo / brand

PPT_ORDER = [
    "Breakfast", "Juice", "Lunch", "Salad", "Soup",
    "Dessert", "Evening Snacks", "Fruit Lunch",
]

_EMBED = qn('r:embed')
_RID   = qn('r:id')


# ── slide helpers ─────────────────────────────────────────────────────────────

def _clone(prs, src_idx):
    """Clone prs.slides[src_idx] and append it; keeps all theme/fonts/media."""
    src  = prs.slides[src_idx]
    tgt  = prs.slides.add_slide(prs.slide_layouts[6])
    tree = tgt.shapes._spTree

    # Wipe default empty children — we replace with exact copy from source
    for ch in list(tree):
        tree.remove(ch)

    # Remap image rIds: source slide → target slide (same image parts, new rIds)
    rmap = {}
    for rId, rel in src.part.rels.items():
        if 'image' in rel.reltype.lower():
            rmap[rId] = tgt.part.relate_to(rel.target_part, rel.reltype)

    for ch in src.shapes._spTree:
        elem = copy.deepcopy(ch)
        for node in elem.iter():
            v = node.get(_EMBED)
            if v in rmap:
                node.set(_EMBED, rmap[v])
        tree.append(elem)

    return tgt


def _delete_slide(prs, idx):
    """Remove slide at idx from the presentation (XML + relationship)."""
    sldIdLst = prs.slides._sldIdLst
    sldId    = sldIdLst[idx]
    rId      = sldId.get(_RID)
    sldIdLst.remove(sldId)
    prs.part.drop_rel(rId)


def _find(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def _find_in(group, name):
    try:
        for s in group.shapes:
            if s.name == name:
                return s
    except Exception:
        pass
    return None


def _text(shape, val, autofit=False):
    """Set text on first run of first paragraph, preserving template font."""
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    while len(tf.paragraphs) > 1:
        tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
    para = tf.paragraphs[0]
    while len(para.runs) > 1:
        para.runs[-1]._r.getparent().remove(para.runs[-1]._r)
    if para.runs:
        para.runs[0].text = val
    else:
        para.add_run().text = val
    if autofit:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE



def _clear_text(shape):
    """Remove all text from a shape's text frame (wipe placeholder notes)."""
    if not shape or not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.text = ''


def _to_supported_image(img_path: str):
    """Return a BytesIO with JPEG data; converts WEBP/AVIF/unsupported formats."""
    import io
    from PIL import Image as _PILImage
    _SUPPORTED = {'JPEG', 'PNG', 'BMP', 'GIF', 'TIFF', 'WMF'}
    img = _PILImage.open(img_path)
    if img.format in _SUPPORTED:
        return img_path          # already fine — pass path directly
    buf = io.BytesIO()
    img.convert('RGB').save(buf, format='JPEG')
    buf.seek(0)
    return buf


def _blip(shape, img_path, slide):
    """Swap the picture fill of a freeform shape with the given image file."""
    if not shape or not img_path or not Path(img_path).exists():
        return
    try:
        src = _to_supported_image(img_path)
        _, rId = slide.part.get_or_add_image_part(src)
        b = shape._element.find('.//' + qn('a:blip'))
        if b is not None:
            b.set(_EMBED, rId)
    except Exception:
        pass


def _clear_blip(shape):
    """Replace picture fill with a neutral solid fill (no image shown)."""
    if not shape:
        return
    try:
        spPr = shape._element.find(qn('p:spPr'))
        if spPr is None:
            return
        blipFill = spPr.find(qn('a:blipFill'))
        if blipFill is not None:
            idx = list(spPr).index(blipFill)
            spPr.remove(blipFill)
            solidFill = _et.Element(qn('a:solidFill'))
            srgbClr = _et.SubElement(solidFill, qn('a:srgbClr'))
            srgbClr.set('val', 'C8E1EB')
            spPr.insert(idx, solidFill)
    except Exception:
        pass


def _build_nutrition(slide, dish):
    """Replace template nutrition boxes with a clean white panel, 5 single-line rows."""
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    OLIVE  = RGBColor(64,  94,  1)
    WHITE  = RGBColor(255, 255, 255)

    # Wipe old template label/value text boxes (they'll be behind the white rect)
    for name in ('TextBox 19', 'TextBox 20', 'TextBox 21', 'TextBox 22',
                 'TextBox 23', 'TextBox 24', 'TextBox 25', 'TextBox 26'):
        s = _find(slide, name)
        if s:
            _clear_text(s)

    # White rectangle — spans below NUTRITION FACTS header, covers all 5 rows
    BOX_L = 6125973          # Table 17 left
    BOX_T = 3425563          # just below NUTRITION FACTS header
    BOX_W = 5100534          # Table 17 width
    N, ROW_H = 5, 530000
    BOX_H = N * ROW_H

    rect = slide.shapes.add_shape(1, Emu(BOX_L), Emu(BOX_T), Emu(BOX_W), Emu(BOX_H))
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = OLIVE
    rect.line.width = Pt(1.5)

    rows = [
        ('CALORIES', dish.get('calories', '—'), 'KCAL'),
        ('PROTEIN',  dish.get('protein',  '—'), 'g'),
        ('FATS',     dish.get('fat',      '—'), 'g'),
        ('FIBRE',    dish.get('fiber',    '—'), 'g'),
        ('CARBS',    dish.get('carbs',    '—'), 'g'),
    ]

    PAD   = 100000
    LBL_W = 2000000
    VAL_L = BOX_L + LBL_W + PAD * 2
    VAL_W = BOX_W - LBL_W - PAD * 3

    for i, (lbl, val, unit) in enumerate(rows):
        ry  = BOX_T + ROW_H * i
        sv  = str(val).strip() if val not in ('—', '', None) else '—'
        vs  = sv if sv == '—' or unit.lower() in sv.lower() else f'{sv} {unit}'

        lb = slide.shapes.add_textbox(Emu(BOX_L + PAD), Emu(ry + PAD // 2),
                                      Emu(LBL_W), Emu(ROW_H - PAD))
        lb.text_frame.word_wrap = False
        lp = lb.text_frame.paragraphs[0]
        lp.alignment = PP_ALIGN.LEFT
        lr = lp.add_run()
        lr.text = lbl
        lr.font.size = Pt(22)
        lr.font.bold = True
        lr.font.color.rgb = OLIVE

        vb = slide.shapes.add_textbox(Emu(VAL_L), Emu(ry + PAD // 2),
                                      Emu(VAL_W), Emu(ROW_H - PAD))
        vb.text_frame.word_wrap = False
        vp = vb.text_frame.paragraphs[0]
        vp.alignment = PP_ALIGN.RIGHT
        vr = vp.add_run()
        vr.text = vs
        vr.font.size = Pt(22)
        vr.font.bold = True
        vr.font.color.rgb = OLIVE


# ── slide builders ────────────────────────────────────────────────────────────

def _good_morning(prs, welcome_text, company_image):
    slide = _clone(prs, _T_GM)
    _text(_find(slide, 'TextBox 11'), welcome_text)
    # Remove the team logo shape (Freeform 12) from the Good Morning slide
    shape = _find(slide, 'Freeform 12')
    if shape:
        el = shape._element
        el.getparent().remove(el)


def _dish_slide(prs, dish, section):
    slide = _clone(prs, _T_DISH)

    _text(_find(slide, 'TextBox 8'),  dish.get('name', '').upper(), autofit=True)
    
    # Category section text size refined to 54 Pt, word wrap disabled to prevent layout overlap
    sec_shape = _find(slide, 'TextBox 27')
    if sec_shape:
        _text(sec_shape, section, autofit=False)
        if sec_shape.has_text_frame:
            sec_shape.text_frame.word_wrap = False
            from pptx.util import Pt
            for para in sec_shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(54)

    # Dish photo — replace blob fill; clear border placeholder text
    grp = _find(slide, 'Group 9')
    if grp:
        _clear_text(_find_in(grp, 'Freeform 11'))
        fr10 = _find_in(grp, 'Freeform 10')
        img  = dish.get('image_path', '')
        from database import resolve_image_path
        resolved = resolve_image_path(img)
        if resolved and Path(resolved).exists():
            _blip(fr10, resolved, slide)
        else:
            _clear_blip(fr10)

    # Nutrition: white box with 5 clean single-line rows
    _build_nutrition(slide, dish)


def _team_photo_slide(prs, team_photo):
    slide = _clone(prs, _T_TEAM)
    grp = _find(slide, 'Group 12')
    if grp:
        fr = _find_in(grp, 'Freeform 13')
        _clear_text(fr)
        if team_photo:
            _blip(fr, team_photo, slide)


def _logo_slide(prs, company_image, company_name):
    slide = _clone(prs, _T_LOGO)
    
    # Team name text size refined to 40 Pt, word wrap disabled to prevent overlap
    shape = _find(slide, 'TextBox 14')
    if shape:
        _text(shape, company_name, autofit=False)
        if shape.has_text_frame:
            shape.text_frame.word_wrap = False
            from pptx.util import Pt
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(40)
                    
    grp = _find(slide, 'Group 12')
    if grp:
        fr = _find_in(grp, 'Freeform 13')
        _clear_text(fr)
        if company_image:
            _blip(fr, company_image, slide)


# ── main export ───────────────────────────────────────────────────────────────

def export_pptx(
    week_data: dict,
    day: str,
    week_start: date,
    dishes_db: dict,
    output_path,
    company_name: str = "Chathradhari Caterers",
    canteen_name:  str = "Oyster",
    welcome_text:  str = "Welcome to Oyster Cafeteria",
    company_image: str = "",
    team_photo:    str = "",
    custom_slides: list | None = None,
) -> Path:
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    from database import resolve_image_path
    company_image = resolve_image_path(company_image)
    team_photo = resolve_image_path(team_photo)

    # Load template — slides 0-3 stay at fixed indices throughout the build
    prs = Presentation(str(TEMPLATE_PATH))

    # Good Morning slide
    _good_morning(prs, welcome_text, company_image)

    # Dish slides grouped by PPT section
    if week_data and day in week_data:
        day_slots = week_data[day]
    elif week_data:
        day_slots = week_data
    else:
        day_slots = {}
    section_dishes: dict[str, list] = {s: [] for s in PPT_ORDER}

    for slot_id, _, slot_sec, ppt_sec, _ in SLOTS:
        dish_names_str = day_slots.get(slot_id, "").strip()
        if not dish_names_str:
            continue
        # Split by '+' to handle multiple dishes in a single slot
        names = [n.strip() for n in dish_names_str.split("+") if n.strip()]
        for dish_name in names:
            dish_info = dict(dishes_db.get(dish_name, {}))
            dish_info.setdefault("name", dish_name)
            
            db_cat = dish_info.get("category", "")
            
            # Apply dynamic section-wise slide labeling rules
            if slot_sec == "FRUIT":
                slide_sec = "Fruit Lunch"
            elif db_cat in ("Dessert", "Sweet"):
                slide_sec = "Dessert"
            elif db_cat == "Juice":
                slide_sec = "Juice"
            elif slot_sec == "LUNCH":
                slide_sec = "Lunch"
            elif slot_id in ("MORN1", "MORN2", "MORN3"):
                slide_sec = "Breakfast"
            elif slot_sec == "ADDITIONAL":
                slide_sec = "Breakfast"
            else:
                slide_sec = ppt_sec

            if slide_sec in section_dishes:
                section_dishes[slide_sec].append(dish_info)

    for ppt_sec in PPT_ORDER:
        for dish_info in section_dishes[ppt_sec]:
            _dish_slide(prs, dish_info, ppt_sec)

    # Closing slides
    _team_photo_slide(prs, team_photo)
    _logo_slide(prs, company_image, company_name)

    # Remove the 4 original template slides (always at indices 0-3)
    for _ in range(4):
        _delete_slide(prs, 0)

    prs.save(str(output_path))
    return output_path
