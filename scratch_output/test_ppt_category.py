import sys
from pathlib import Path
sys.path.append('menu_planner')
import database as db
from export.pptx_gen import export_pptx

def test():
    db.seed_database_on_first_run()
    db.init_db()
    
    # We know plan_id 4 has Omelet in EGG slot on TUESDAY
    # Let's get week data for plan_id 4
    week_data = db.get_week_data(4)
    print("Week data for plan 4:", week_data)
    
    dishes_db = {d["name"]: d for d in db.get_all_dishes()}
    
    out_path = Path("scratch_output/test_plan4.pptx")
    export_pptx(
        week_data=week_data,
        day="TUESDAY",
        week_start=db.get_all_plans()[0]["week_start"], # dummy, not used for slide categorisation anyway
        dishes_db=dishes_db,
        output_path=out_path
    )
    print("PPT generated successfully at", out_path)
    
    # Let's inspect the shapes on each slide
    from pptx import Presentation
    prs = Presentation(str(out_path))
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i} ---")
        for s in slide.shapes:
            if s.has_text_frame:
                text = s.text_frame.text.strip()
                if text:
                    print(f"Shape: {s.name} | Text: {text}")

if __name__ == "__main__":
    test()
