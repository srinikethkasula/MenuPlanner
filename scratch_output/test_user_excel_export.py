import sys
import os
import openpyxl
from pathlib import Path
from datetime import date

sys.path.append('menu_planner')
import database as db
from export.pptx_gen import export_pptx

def import_and_test():
    db.seed_database_on_first_run()
    db.init_db()
    
    # Let's load the excel file
    excel_path = r"C:\Users\kasul\Downloads\test folder 1\Weekly Menu 15 to 19 June 26.xlsx"
    if not os.path.exists(excel_path):
        print("Excel file not found!")
        return
        
    wb = openpyxl.load_workbook(excel_path)
    ws = wb.active
    
    # Let's read the dates
    # Row 6 has: DATE, 15-Jun-26, 16-Jun-26, 17-Jun-26, 18-Jun-26, 19-Jun-26
    # Row 5 has headers: LIST OF ITEMS, MONDAY, TUESDAY, etc.
    # We will map each cell value to slot assignments
    from config import SLOTS, DAYS
    
    # We will create a plan for 2026-06-15
    week_start = date(2026, 6, 15)
    plan_id = db.get_or_create_plan(week_start)
    
    # Read slots and assign
    # In Excel:
    # Row 7 is LUNCH header
    # Row 8 is NON.VEG
    # Row 9 is SPL.VEG.
    # ...
    # Let's match display names or row names
    # Column 3 has the slot name, Columns 4 to 8 have Monday to Friday values
    for row in range(7, 45):
        slot_name_val = ws.cell(row=row, column=3).value
        if not slot_name_val:
            continue
        slot_name = str(slot_name_val).strip()
        
        # Let's find the slot_id matching this slot_name (display name)
        # Or match by position:
        # Row 8: NON_VEG
        # Row 9: SPL_VEG
        # etc.
        # Let's match slot_id by iterating over SLOTS and checking display names
        slot_id = None
        for s_id, disp, sec, ppt, _ in SLOTS:
            if disp.strip().lower() == slot_name.lower():
                slot_id = s_id
                break
        
        if not slot_id:
            # Try matching display names with trailing space or slightly different
            for s_id, disp, sec, ppt, _ in SLOTS:
                if slot_name.lower() in disp.lower() or disp.lower() in slot_name.lower():
                    slot_id = s_id
                    break
                    
        if not slot_id:
            continue
            
        for d_idx, day in enumerate(DAYS):
            col = 4 + d_idx
            dish_name = ws.cell(row=row, column=col).value
            if dish_name:
                dish_name = str(dish_name).strip()
                db.set_slot(plan_id, day, slot_id, dish_name)
                
    print(f"Imported assignments for plan_id {plan_id} (week starting {week_start})")
    
    # Now let's export PPT for each day and print the results
    dishes_db = {d["name"]: d for d in db.get_all_dishes()}
    
    for day in DAYS:
        out_path = Path(f"scratch_output/test_plan5_{day.lower()}.pptx")
        export_pptx(
            week_data=db.get_week_data(plan_id),
            day=day,
            week_start=week_start,
            dishes_db=dishes_db,
            output_path=out_path
        )
        print(f"\n--- PPT for {day} ---")
        from pptx import Presentation
        prs = Presentation(str(out_path))
        for i, slide in enumerate(prs.slides):
            dish_names = [s.text_frame.text.strip() for s in slide.shapes if s.name == "TextBox 8"]
            categories = [s.text_frame.text.strip() for s in slide.shapes if s.name == "TextBox 27"]
            if dish_names:
                print(f"Slide {i}: Dish = {dish_names} | Category = {categories}")

if __name__ == "__main__":
    import_and_test()
